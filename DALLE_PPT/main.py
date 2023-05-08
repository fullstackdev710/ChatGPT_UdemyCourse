import tkinter as tk
from pptx import Presentation
from pptx.util import Inches, Pt
import openai
from io import BytesIO
from pathlib import Path
import requests

openai.api_key = "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"


def slide_generator(text, prs):
    prompt = f"Summarize the following text to a DALL-E image generation" \
        f"prompt: \n {text}"

    model_engine = "text-davinci-003"
    dlp = openai.Completion.create(
        engine=model_engine,
        prompt=prompt,
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )

    dalle_prompt = dlp.choices[0].text
    response = openai.Image.create(
        prompt=dalle_prompt+"Style: digital art",
        n=1,
        size="1024x1024"
    )
    image_url = response["data"][0]['url']

    prompt = f"Create a bullet point text for a Powerpoint" \
        f"slide from the following text: \n {text}"
    ppt = openai.Completion.create(
        engine=model_engine,
        prompt=prompt,
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_text = ppt.choices[0].text

    prompt = f"Create a title for a Powerpoint" \
        f"slide from the following text: \n {text}"
    ppt = openai.Completion.create(
        engine=model_engine,
        prompt=prompt,
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_header = ppt.choices[0].text

    # Add a new slide to the presetation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    response = requests.get(image_url)
    img_bytes = BytesIO(response.content)
    slide.shapes.add_picture(img_bytes, Inches(1), Inches(1))

    # Add text box
    txtBox = slide.shapes.add_textbox(
        Inches(3), Inches(1), Inches(4), Inches(1.5))
    tf = txtBox.text_frame
    tf.text = ppt_text
    title_shape = slide.shapes.title
    title_shape.text = ppt_header


def get_slides():
    """
    This is paragraph 1

    This is a paragraph 2
    ["This is paragraph 1", "This is a paragraph 2"]
    """
    text = text_field.get("1.0", "end-1c")
    paragraphs = text.split("\n\n")
    prs = Presentation()
    width = Pt(1920)
    height = Pt(1080)
    prs.slide_width = width
    prs.slide_height = height
    for paragraph in paragraphs:
        slide_generator(paragraph, prs)

    prs.save("my_presentation.pptx")


app = tk.Tk()
app.title("Create PPT Slides")
app.geometry("800x600")

# Create text field
text_field = tk.Text(app)
text_field.pack(fill="both", expand=True)
text_field.configure(wrap="word", font=("Arial", 12))
text_field.focus_set()

# Create the button to create slides
create_button = tk.Button(app, text="Create Slides", command=get_slides())
create_button.pack()

app.mainloop()
